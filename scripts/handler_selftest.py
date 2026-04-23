"""Smoke-test installed file-type handlers against minimal samples.

This module is invoked from `/vault-bridge:setup` Step 6.5f as an
acceptance gate: for each installed handler, generate or embed a tiny
sample file of the matching type, invoke the stub's `read_text` and /
or `extract_images` as claimed by its CAPABILITIES, and assert the
output is non-empty and the claimed kind.

The intent is to catch regressions where a handler was generated from
the generic stub template (`# TODO: implement`) and advertises
capabilities it does not actually deliver.

Public API
----------
    SelftestResult       — dataclass with per-extension outcome
    run_selftest(...)    — runs every .py in `handlers_dir` against a sample
    generate_sample(ext) — returns (bytes, suffix) or (None, None) when
                           no sample generator exists for `ext`

Never raises. All import / decode / exec failures are captured in the
SelftestResult.error field. Extensions with no generator are reported
as `skipped=True`.

Python 3.9 compatible.
"""
from __future__ import annotations

import importlib.util
import io
import logging
import shutil
import struct
import tempfile
import zlib
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)


# Extensions that cannot be smoke-tested with a synthetic sample —
# generating them requires either an external tool we can't depend on
# or a binary fixture we don't ship. Each value is the reason string
# reported back to the user so setup output is honest about WHY the
# smoke test was skipped.
_SKIP_REASONS: Dict[str, str] = {
    "dwg": "requires ODA File Converter (external tool)",
    "3dm": "no synthetic generator (rhino3dm cannot write from Python)",
    "doc": "no synthetic generator (legacy OLE binary not writable from Python)",
    "ppt": "no synthetic generator (legacy OLE binary not writable from Python)",
    "psd": "no synthetic generator (needs real PSD fixture)",
}


# ---------------------------------------------------------------------------
# Dataclass
# ---------------------------------------------------------------------------

@dataclass
class SelftestResult:
    """Outcome of a single handler selftest."""
    ext: str
    handler_path: str
    ok: bool = False
    skipped: bool = False
    skip_reason: str = ""
    read_text_ran: bool = False
    read_text_non_empty: bool = False
    extract_images_ran: bool = False
    extract_images_non_empty: bool = False
    error: str = ""


# ---------------------------------------------------------------------------
# Sample byte constants
# ---------------------------------------------------------------------------

# Minimal valid single-pixel PNG (red).
_PNG_SAMPLE: bytes = (
    b"\x89PNG\r\n\x1a\n"
    b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
    b"\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0\x00\x00\x00\x03\x00\x01\x5e\xf3\x2a\x0a"
    b"\x00\x00\x00\x00IEND\xaeB`\x82"
)

# Minimal valid single-pixel JPEG (black).
_JPG_SAMPLE: bytes = bytes.fromhex(
    "ffd8ffdb004300080606070605080707070909080a0c140d0c0b0b0c1912130f141d1a"
    "1f1e1d1a1c1c20242e2720222c231c1c2837292c30313434341f27393d38323c2e3332"
    "ffc0000b080001000101011100ffc4001f0000010501010101010100000000000000"
    "000102030405060708090a0bffc400b5100002010303020403050504040000017d01"
    "02030004110512213141061351610722718132081442912a1b1c1d1e242526272829"
    "2a333435363738393a434445464748494a535455565758595a636465666768696a73"
    "7475767778797a838485868788898a92939495969798999aa2a3a4a5a6a7a8a9aab2"
    "b3b4b5b6b7b8b9bac2c3c4c5c6c7c8c9cad2d3d4d5d6d7d8d9dae1e2e3e4e5e6e7e8"
    "e9eaf1f2f3f4f5f6f7f8f9faffda0008010100003f00fbd0ffd9"
)

# Minimal 2-byte GIF header + trailer (renders as 1x1 black in most decoders).
_GIF_SAMPLE: bytes = (
    b"GIF89a\x01\x00\x01\x00\x80\x00\x00\x00\x00\x00\xff\xff\xff,\x00\x00\x00\x00"
    b"\x01\x00\x01\x00\x00\x02\x02\x44\x01\x00;"
)

# Minimal 4-byte BMP (2x2 solid white, 24-bit).
_BMP_SAMPLE: bytes = bytes.fromhex(
    "424d3a000000000000003600000028000000010000000100000001001800000000000"
    "400000000000000000000000000000000000000ffffff00"
)

# Minimal valid single-pixel RGB WEBP (Lossy VP8).
_WEBP_SAMPLE: bytes = bytes.fromhex(
    "524946461a000000574542505650384c0d0000002f00000000006f0100000000"
)


def _make_rtf() -> bytes:
    return b"{\\rtf1\\ansi Hello\\par}"


def _make_txt() -> bytes:
    return "Hello from vault-bridge selftest — 你好 🌱".encode("utf-8")


def _make_md() -> bytes:
    return b"# vault-bridge\n\nSelftest markdown sample.\n"


def _make_pdf() -> bytes:
    """Minimal single-page PDF with one line of text."""
    # Hand-rolled so we don't pull reportlab as a dep just for selftests.
    body = (
        b"%PDF-1.4\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Count 1/Kids[3 0 R]>>endobj\n"
        b"3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 200 100]/Contents 4 0 R"
        b"/Resources<</Font<</F1 5 0 R>>>>>>endobj\n"
        b"4 0 obj<</Length 44>>stream\n"
        b"BT /F1 12 Tf 10 50 Td (vault-bridge ok) Tj ET\n"
        b"endstream endobj\n"
        b"5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
        b"xref\n0 6\n0000000000 65535 f \n"
        b"0000000009 00000 n \n"
        b"0000000058 00000 n \n"
        b"0000000109 00000 n \n"
        b"0000000200 00000 n \n"
        b"0000000290 00000 n \n"
        b"trailer<</Size 6/Root 1 0 R>>\nstartxref\n349\n%%EOF\n"
    )
    return body


def _make_docx() -> Optional[bytes]:
    """Tiny .docx with a paragraph AND an embedded PNG.

    The handler claims `extract_images=True` so the fixture must
    include at least one image for the smoke-test to exercise that
    capability. Without the embed the selftest false-FAILs even
    though the handler works correctly on real files.
    """
    try:
        import docx  # type: ignore
    except Exception:
        return None
    doc = docx.Document()
    doc.add_paragraph("vault-bridge selftest — docx content line.")
    try:
        doc.add_picture(io.BytesIO(_PNG_SAMPLE))
    except Exception as exc:
        logger.debug("_make_docx: embedded-image insert failed: %s", exc)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx() -> Optional[bytes]:
    """Tiny .pptx with a title AND an embedded PNG (see _make_docx note)."""
    try:
        import pptx  # type: ignore
        from pptx.util import Emu  # type: ignore
    except Exception:
        return None
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    if title is not None:
        title.text = "vault-bridge selftest"
    try:
        slide.shapes.add_picture(
            io.BytesIO(_PNG_SAMPLE),
            left=Emu(914400),
            top=Emu(914400),
            width=Emu(914400),
            height=Emu(914400),
        )
    except Exception as exc:
        logger.debug("_make_pptx: embedded-image insert failed: %s", exc)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx() -> Optional[bytes]:
    """Tiny .xlsx with cells AND an embedded PNG (see _make_docx note)."""
    try:
        import openpyxl  # type: ignore
    except Exception:
        return None
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "vault-bridge"
    ws["B1"] = "selftest"
    try:
        from openpyxl.drawing.image import Image as XLImage  # type: ignore
        # openpyxl.Image wants a file path; write the PNG to a
        # tempfile long enough for wb.save() to consume it.
        with tempfile.NamedTemporaryFile(
            suffix=".png", delete=False
        ) as tf:
            tf.write(_PNG_SAMPLE)
            tf_path = tf.name
        try:
            ws.add_image(XLImage(tf_path), "D4")
            buf = io.BytesIO()
            wb.save(buf)
            return buf.getvalue()
        finally:
            try:
                Path(tf_path).unlink()
            except Exception:
                pass
    except Exception as exc:
        logger.debug("_make_xlsx: embedded-image insert failed: %s", exc)
        # Fall back to an image-less xlsx — better than no sample at all.
        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()


def _make_tiff() -> Optional[bytes]:
    try:
        from PIL import Image  # type: ignore
    except Exception:
        return None
    buf = io.BytesIO()
    Image.new("RGB", (1, 1), color=(255, 0, 0)).save(buf, format="TIFF")
    return buf.getvalue()


def _make_heif() -> Optional[bytes]:
    """Generate a single-pixel HEIF. Requires pillow-heif to be installed."""
    try:
        import pillow_heif  # type: ignore
        try:
            pillow_heif.register_heif_opener()
        except Exception:
            pass
        from PIL import Image  # type: ignore
    except Exception:
        return None
    buf = io.BytesIO()
    try:
        Image.new("RGB", (1, 1), color=(0, 255, 0)).save(
            buf, format="HEIF"
        )
    except Exception as exc:
        logger.debug("_make_heif: save failed: %s", exc)
        return None
    return buf.getvalue()


def _make_dxf() -> Optional[bytes]:
    """Generate a tiny DXF with one TEXT entity via ezdxf."""
    try:
        import ezdxf  # type: ignore
    except Exception:
        return None
    try:
        doc = ezdxf.new("R2010")
        msp = doc.modelspace()
        # ezdxf 1.x: `add_text` returns a Text entity; the preferred
        # placement API differs across releases. Passing dxfattribs
        # avoids the `.set_pos`/`.set_placement` split.
        msp.add_text("vault-bridge", dxfattribs={"insert": (0, 0)})
        buf = io.StringIO()
        doc.write(buf)
        return buf.getvalue().encode("utf-8")
    except Exception as exc:
        logger.debug("_make_dxf: failed: %s", exc)
        return None


def _make_ai() -> bytes:
    """Modern .ai is a PDF container — reuse the PDF sample."""
    return _make_pdf()


_SAMPLES: Dict[str, Tuple[bytes, str]] = {}


def _load_samples() -> None:
    """Populate the sample table lazily so dynamic generators run at call time."""
    if _SAMPLES:
        return

    static_map: Dict[str, bytes] = {
        "png": _PNG_SAMPLE,
        "jpg": _JPG_SAMPLE,
        "jpeg": _JPG_SAMPLE,
        "gif": _GIF_SAMPLE,
        "bmp": _BMP_SAMPLE,
        "webp": _WEBP_SAMPLE,
        "txt": _make_txt(),
        "md": _make_md(),
        "rtf": _make_rtf(),
        "pdf": _make_pdf(),
        "ai": _make_ai(),
    }
    for ext, data in static_map.items():
        _SAMPLES[ext] = (data, ext)

    dynamic: Dict[str, Optional[bytes]] = {
        "docx": _make_docx(),
        "pptx": _make_pptx(),
        "xlsx": _make_xlsx(),
        "tiff": _make_tiff(),
        "tif": _make_tiff(),
        "heic": _make_heif(),
        "heif": _make_heif(),
        "dxf": _make_dxf(),
    }
    for ext, data in dynamic.items():
        if data:
            _SAMPLES[ext] = (data, ext)


def generate_sample(ext: str) -> Tuple[Optional[bytes], Optional[str]]:
    """Return (bytes, suffix) for a smoke-test sample, or (None, None)."""
    _load_samples()
    ext_clean = ext.lower().lstrip(".")
    entry = _SAMPLES.get(ext_clean)
    if entry is None:
        return None, None
    return entry


# ---------------------------------------------------------------------------
# Runner
# ---------------------------------------------------------------------------

def _load_handler(handler_path: Path):
    spec = importlib.util.spec_from_file_location(handler_path.stem, str(handler_path))
    if not spec or not spec.loader:
        return None
    mod = importlib.util.module_from_spec(spec)
    try:
        spec.loader.exec_module(mod)  # type: ignore[union-attr]
    except Exception as exc:
        logger.debug("handler load failed for %s: %s", handler_path, exc)
        return None
    return mod


def _ext_from_handler_filename(handler_path: Path) -> str:
    """Derive extension from `<category>_<ext>.py` filename convention."""
    stem = handler_path.stem
    if "_" in stem:
        return stem.rsplit("_", 1)[1]
    return stem


def _selftest_one(handler_path: Path) -> SelftestResult:
    ext = _ext_from_handler_filename(handler_path)
    result = SelftestResult(ext=ext, handler_path=str(handler_path))

    data, suffix = generate_sample(ext)
    if data is None:
        result.skipped = True
        # Prefer the curated reason over a generic "no generator" so
        # setup output distinguishes "can't write .3dm from Python"
        # from "requires ODA File Converter" from "pillow-heif not
        # installed — a real dependency issue that should be fixed".
        if ext in _SKIP_REASONS:
            result.skip_reason = _SKIP_REASONS[ext]
        elif ext in ("heic", "heif"):
            result.skip_reason = "pillow-heif not installed"
        elif ext in ("tif", "tiff"):
            result.skip_reason = "Pillow not installed"
        elif ext == "dxf":
            result.skip_reason = "ezdxf not installed"
        else:
            result.skip_reason = f"no sample generator for .{ext}"
        return result

    mod = _load_handler(handler_path)
    if mod is None:
        result.error = "handler failed to import"
        return result

    caps = getattr(mod, "CAPABILITIES", None)
    if not isinstance(caps, dict):
        result.error = "handler has no CAPABILITIES dict"
        return result

    with tempfile.TemporaryDirectory() as td:
        tmp_path = Path(td) / f"sample.{suffix}"
        tmp_path.write_bytes(data)
        out_dir = Path(td) / "out"
        out_dir.mkdir()

        rt_ok = True
        if caps.get("read_text"):
            result.read_text_ran = True
            try:
                text = mod.read_text(str(tmp_path))
            except Exception as exc:
                result.error = f"read_text raised: {exc}"
                return result
            if isinstance(text, str) and text.strip():
                result.read_text_non_empty = True
            else:
                rt_ok = False

        ei_ok = True
        if caps.get("extract_images"):
            result.extract_images_ran = True
            try:
                images = mod.extract_images(str(tmp_path), str(out_dir))
            except Exception as exc:
                result.error = f"extract_images raised: {exc}"
                return result
            if isinstance(images, list) and len(images) > 0:
                result.extract_images_non_empty = True
            else:
                ei_ok = False

        # A handler passes only when every capability it claims produces
        # real output. A handler that claims neither capability is a stub
        # and fails the gate.
        claimed_any = bool(caps.get("read_text") or caps.get("extract_images"))
        result.ok = claimed_any and rt_ok and ei_ok
    return result


def run_selftest(handlers_dir: Path) -> List[SelftestResult]:
    """Run selftest for every .py handler file in `handlers_dir`.

    Returns a list of SelftestResult, one per handler module. Never raises.
    """
    handlers_dir = Path(handlers_dir)
    results: List[SelftestResult] = []
    if not handlers_dir.exists():
        return results
    for handler_path in sorted(handlers_dir.glob("*.py")):
        if handler_path.name.startswith("_"):
            continue
        try:
            results.append(_selftest_one(handler_path))
        except Exception as exc:
            logger.debug("selftest one failed for %s: %s", handler_path, exc)
            results.append(SelftestResult(
                ext=_ext_from_handler_filename(handler_path),
                handler_path=str(handler_path),
                error=str(exc),
            ))
    return results


def format_summary(results: List[SelftestResult]) -> str:
    """Format a markdown table summary suitable for the setup report."""
    lines = [
        "| ext | status | read_text | extract_images | notes |",
        "| --- | --- | --- | --- | --- |",
    ]
    for r in sorted(results, key=lambda x: x.ext):
        if r.skipped:
            status = "— skipped"
            note = r.skip_reason
        elif r.ok:
            status = "ok"
            note = ""
        else:
            status = "FAIL"
            note = r.error or "capability claimed but output empty"
        rt = "—"
        if r.read_text_ran:
            rt = "ok" if r.read_text_non_empty else "EMPTY"
        ei = "—"
        if r.extract_images_ran:
            ei = "ok" if r.extract_images_non_empty else "EMPTY"
        lines.append(f"| {r.ext} | {status} | {rt} | {ei} | {note} |")
    return "\n".join(lines)
