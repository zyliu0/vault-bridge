"""Generator for scripts/file_type_handlers.py.

Reads file_type_config from <workdir>/.vault-bridge/config.json and writes a
fresh file_type_handlers.py that scan commands can import directly.

The generated script is self-contained — it does NOT import this module.

Public API
----------
    generate(workdir, out_path=None) -> Path
    generate_from_dict(file_type_config, out_path) -> Path

Config shape (file_type_config key in config.json)
---------------------------------------------------
    {
        "category_overrides": {
            "document-pdf": {
                "extract_text": false,
                "extract_images": true,
                "compress": true,
                "run_vision": false
            },
            ...
        },
        "extra_extensions": {
            "document-office": ["pages", "numbers", "key"],
            ...
        },
        "skip_extensions": ["rar", "7z", ...]
    }

All keys are optional. Unknown category names in category_overrides are
silently ignored. The defaults match what the base file_type_handlers.py
defines.

Python 3.9 compatible.
"""
import json
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Default category definitions (mirrors file_type_handlers.py)
# ---------------------------------------------------------------------------

@dataclass
class _CategoryDef:
    extract_text: bool
    extract_images: bool
    compress: bool
    run_vision: bool


_DEFAULT_CATEGORIES: Dict[str, _CategoryDef] = {
    "document-pdf":    _CategoryDef(extract_text=True,  extract_images=True,  compress=True,  run_vision=True),
    "document-office": _CategoryDef(extract_text=True,  extract_images=True,  compress=True,  run_vision=False),
    "image-raster":    _CategoryDef(extract_text=False, extract_images=True,  compress=True,  run_vision=True),
    "image-vector":    _CategoryDef(extract_text=False, extract_images=True,  compress=False, run_vision=False),
    "video":           _CategoryDef(extract_text=False, extract_images=False, compress=False, run_vision=False),
    "audio":           _CategoryDef(extract_text=False, extract_images=False, compress=False, run_vision=False),
    "text-plain":      _CategoryDef(extract_text=True,  extract_images=False, compress=False, run_vision=False),
    "archive":         _CategoryDef(extract_text=False, extract_images=False, compress=False, run_vision=False),
}

# Extensions explicitly skipped by default. Users can still re-add these via
# `file_type_config.extra_extensions` if they want handling, but by default
# we do NOT generate handlers for render passes, mesh exchange formats, HDR
# images, Grasshopper scripts, JSON config, or Rhino backups.
_DEFAULT_SKIP_EXTENSIONS: List[str] = [
    "exr", "hdr",       # HDR / OpenEXR render passes — no readable content
    "gh", "ghx",        # Rhino Grasshopper definitions — binary scripts
    "obj", "stl",       # 3D mesh exchange — geometry only, no narrative
    "json",             # machine-readable config; not diary-worthy
    "3dmbak",           # Rhino backup files; alias redundant with .3dm
]

# Default extensions per category (mirrors file_type_handlers.py HANDLERS)
_DEFAULT_EXTENSIONS: Dict[str, List[str]] = {
    "document-pdf":    ["pdf"],
    "document-office": ["docx", "doc", "pptx", "ppt", "xlsx", "xls"],
    "image-raster":    ["jpg", "jpeg", "png", "gif", "webp", "bmp", "tiff", "tif", "heic", "heif"],
    "image-vector":    ["svg", "ai", "eps"],
    "video":           ["mp4", "mov", "avi", "mkv", "wmv", "m4v", "webm"],
    "audio":           ["mp3", "wav", "flac", "aac", "ogg", "m4a"],
    "text-plain":      ["txt", "md", "rst", "csv"],
    "archive":         ["zip", "rar", "7z", "tar"],
}


# ---------------------------------------------------------------------------
# Config helpers
# ---------------------------------------------------------------------------

def _load_file_type_config(workdir: Path) -> Dict[str, Any]:
    """Read file_type_config from config.json; returns {} on any failure."""
    cfg_path = Path(workdir) / ".vault-bridge" / "config.json"
    if not cfg_path.exists():
        return {}
    try:
        data = json.loads(cfg_path.read_text(encoding="utf-8"))
        return dict(data.get("file_type_config") or {})
    except Exception as exc:
        logger.debug("Could not read file_type_config from %s: %s", cfg_path, exc)
        return {}


def _apply_config(
    file_type_config: Dict[str, Any],
) -> tuple:
    """Return (categories_dict, ext_map, installed_readers) after applying overrides.

    Returns:
        categories:        dict[category_slug, _CategoryDef]
        ext_map:           dict[ext, category_slug]  (extension -> category)
        installed_readers: dict[ext, module_path]    (from installed_packages key)
    """
    import copy

    categories = {k: copy.copy(v) for k, v in _DEFAULT_CATEGORIES.items()}
    ext_map: Dict[str, str] = {}

    # Build ext_map from defaults first
    for cat, exts in _DEFAULT_EXTENSIONS.items():
        for ext in exts:
            ext_map[ext] = cat

    # 1. category_overrides — modify category flags
    for cat_name, overrides in (file_type_config.get("category_overrides") or {}).items():
        if cat_name not in categories:
            logger.debug("Unknown category_override '%s' — ignored", cat_name)
            continue
        cat_def = categories[cat_name]
        for flag in ("extract_text", "extract_images", "compress", "run_vision"):
            if flag in overrides:
                setattr(cat_def, flag, bool(overrides[flag]))

    # 2. extra_extensions — add new extensions to categories
    for cat_name, new_exts in (file_type_config.get("extra_extensions") or {}).items():
        if cat_name not in categories:
            logger.debug("Unknown category for extra_extensions '%s' — ignored", cat_name)
            continue
        for ext in new_exts:
            ext_clean = str(ext).lower().lstrip(".")
            if ext_clean:
                ext_map[ext_clean] = cat_name

    # 3. skip_extensions — remove extensions entirely. Merges built-in
    # defaults with any user-supplied skips from config.
    skip_list = list(_DEFAULT_SKIP_EXTENSIONS)
    skip_list.extend(file_type_config.get("skip_extensions") or [])
    for ext in skip_list:
        ext_clean = str(ext).lower().lstrip(".")
        ext_map.pop(ext_clean, None)

    # 4. installed_packages — ext -> module path for dynamic dispatch
    installed_readers: Dict[str, str] = {}
    raw_installed = file_type_config.get("installed_packages") or {}
    if isinstance(raw_installed, dict):
        for ext, mod_path in raw_installed.items():
            ext_clean = str(ext).lower().lstrip(".")
            if ext_clean and mod_path:
                installed_readers[ext_clean] = str(mod_path)

    return categories, ext_map, installed_readers


# ---------------------------------------------------------------------------
# Code generation
# ---------------------------------------------------------------------------

def _render_source(
    categories: Dict[str, _CategoryDef],
    ext_map: Dict[str, str],
    installed_readers: Optional[Dict[str, str]] = None,
) -> str:
    """Render the full Python source for file_type_handlers.py.

    Args:
        categories:        Category definitions (name -> _CategoryDef).
        ext_map:           Extension to category mapping.
        installed_readers: Optional mapping of ext -> handler module path.
                           When provided, read_text/extract_images will dispatch
                           to the named module via importlib.import_module.
    """
    if installed_readers is None:
        installed_readers = {}

    # Build per-category variable name
    def _varname(cat: str) -> str:
        return "_" + cat.upper().replace("-", "_")

    # Render HandlerConfig instances for each category
    cat_defs_lines: List[str] = []
    for cat, defn in categories.items():
        vn = _varname(cat)
        cat_defs_lines.append(
            f"{vn} = HandlerConfig(\n"
            f"    category={cat!r},\n"
            f"    extract_text={defn.extract_text!r},\n"
            f"    extract_images={defn.extract_images!r},\n"
            f"    compress={defn.compress!r},\n"
            f"    run_vision={defn.run_vision!r},\n"
            f")"
        )
    cat_defs_block = "\n".join(cat_defs_lines)

    # Render HANDLERS dict entries
    handlers_lines: List[str] = []
    for ext in sorted(ext_map):
        cat = ext_map[ext]
        vn = _varname(cat)
        handlers_lines.append(f"    {ext!r}: {vn},")
    handlers_block = "\n".join(handlers_lines)

    # Build _INSTALLED_READERS dict literal
    installed_readers_lines: List[str] = []
    for ext_key in sorted(installed_readers):
        mod_path = installed_readers[ext_key]
        installed_readers_lines.append(f"    {ext_key!r}: {mod_path!r},")
    if installed_readers_lines:
        installed_readers_block = (
            "_INSTALLED_READERS: Dict[str, str] = {\n"
            + "\n".join(installed_readers_lines)
            + "\n}"
        )
    else:
        installed_readers_block = "_INSTALLED_READERS: Dict[str, str] = {}"

    # Build the source as a list of lines — avoids triple-quoted f-string
    # indentation problems with textwrap.dedent.
    lines: List[str] = [
        "# AUTO-GENERATED by scripts/generate_file_type_handlers.py",
        "# Generated by vault-bridge setup. Do not edit manually.",
        "# To regenerate: run /vault-bridge:setup or call generate_file_type_handlers.generate()",
        '"""File-type handler registry for vault-bridge.',
        "",
        "Auto-generated — do not edit manually.",
        "To regenerate after changing config, run /vault-bridge:setup or:",
        "    python scripts/generate_file_type_handlers.py",
        '"""',
        "import importlib.util",
        "import logging",
        "from dataclasses import dataclass, field",
        "from pathlib import Path",
        "from typing import Dict, List, Optional",
        "",
        "logger = logging.getLogger(__name__)",
        "",
        "",
        "# ---------------------------------------------------------------------------",
        "# HandlerConfig",
        "# ---------------------------------------------------------------------------",
        "",
        "@dataclass(frozen=True)",
        "class HandlerConfig:",
        "    category: str",
        "    extract_text: bool",
        "    extract_images: bool",
        "    compress: bool",
        "    run_vision: bool",
        "",
        "",
        "# ---------------------------------------------------------------------------",
        "# Category descriptors (auto-generated from config)",
        "# ---------------------------------------------------------------------------",
        "",
        cat_defs_block,
        "",
        "",
        "# ---------------------------------------------------------------------------",
        "# HANDLERS registry",
        "# ---------------------------------------------------------------------------",
        "",
        "HANDLERS: Dict[str, HandlerConfig] = {",
        handlers_block,
        "}",
        "",
        "",
        "# ---------------------------------------------------------------------------",
        "# _INSTALLED_READERS — ext -> handler module path (managed by setup)",
        "# ---------------------------------------------------------------------------",
        "",
        installed_readers_block,
        "",
        "",
        "# ---------------------------------------------------------------------------",
        "# HandlerResult",
        "# ---------------------------------------------------------------------------",
        "",
        "@dataclass",
        "class HandlerResult:",
        "    text: Optional[str] = None",
        "    images: List[Path] = field(default_factory=list)",
        "    skipped: bool = False",
        "    category: str = ''",
        "",
        "",
        "# ---------------------------------------------------------------------------",
        "# Public API",
        "# ---------------------------------------------------------------------------",
        "",
        "def get_handler(path: str) -> Optional[HandlerConfig]:",
        "    if not path:",
        "        return None",
        "    from pathlib import Path as _Path",
        "    p = _Path(path)",
        "    suffix = p.suffix",
        "    if not suffix:",
        "        return None",
        "    ext = suffix.lstrip('.').lower()",
        "    if not ext:",
        "        return None",
        "    return HANDLERS.get(ext)",
        "",
        "",
        "# ---------------------------------------------------------------------------",
        "# Installed-handler loader",
        "# ---------------------------------------------------------------------------",
        "",
        "# Resolve handler stubs relative to the workdir root. This generated file",
        "# lives at <workdir>/scripts/file_type_handlers.py, so parent.parent is",
        "# the workdir. Handler stubs live under <workdir>/.vault-bridge/handlers/.",
        "try:",
        "    _WORKDIR_ROOT = Path(__file__).resolve().parent.parent",
        "except NameError:",
        "    # exec()'d without __file__ — fall back to cwd (test harness only).",
        "    _WORKDIR_ROOT = Path.cwd()",
        "_HANDLERS_SUBDIR = _WORKDIR_ROOT / '.vault-bridge' / 'handlers'",
        "",
        "_HANDLER_MOD_CACHE: Dict[str, object] = {}",
        "",
        "",
        "def _resolve_handler_path(raw: str) -> Optional[Path]:",
        "    \"\"\"Resolve a stored handler identifier to an absolute file path.",
        "",
        "    Accepts all of: bare stub filename ('cad_3dm_3dm.py'),",
        "    workdir-relative path ('.vault-bridge/handlers/cad_3dm_3dm.py'),",
        "    or legacy dotted module name ('handlers.cad_3dm_3dm').",
        "    \"\"\"",
        "    if not raw:",
        "        return None",
        "    # Dotted module name (no slash, no .py suffix) → extract last segment.",
        "    if '/' not in raw and not raw.endswith('.py'):",
        "        leaf = raw.rsplit('.', 1)[-1] if '.' in raw else raw",
        "        candidate = _HANDLERS_SUBDIR / f'{leaf}.py'",
        "        return candidate if candidate.exists() else None",
        "    # Path form: absolute, workdir-relative, or bare filename.",
        "    clean = raw.lstrip('./')",
        "    p = Path(clean)",
        "    if p.is_absolute():",
        "        return p if p.exists() else None",
        "    if clean.startswith('.vault-bridge/'):",
        "        candidate = _WORKDIR_ROOT / clean",
        "    else:",
        "        candidate = _HANDLERS_SUBDIR / Path(clean).name",
        "    return candidate if candidate.exists() else None",
        "",
        "",
        "def _load_installed(ext: str):",
        "    \"\"\"Load the installed handler module for an extension. Cached per ext.\"\"\"",
        "    if ext in _HANDLER_MOD_CACHE:",
        "        return _HANDLER_MOD_CACHE[ext]",
        "    raw = _INSTALLED_READERS.get(ext)",
        "    abs_path = _resolve_handler_path(raw) if raw else None",
        "    if abs_path is None:",
        "        if raw:",
        "            logger.debug('Installed handler for %s not found (stored: %r)', ext, raw)",
        "        _HANDLER_MOD_CACHE[ext] = None",
        "        return None",
        "    try:",
        "        mod_name = f'vb_handler_{ext}_{abs_path.stem}'",
        "        spec = importlib.util.spec_from_file_location(mod_name, abs_path)",
        "        if spec is None or spec.loader is None:",
        "            raise ImportError(f'could not build spec for {abs_path}')",
        "        mod = importlib.util.module_from_spec(spec)",
        "        spec.loader.exec_module(mod)",
        "        _HANDLER_MOD_CACHE[ext] = mod",
        "        return mod",
        "    except Exception as exc:",
        "        logger.debug('Failed to load handler %s: %s', abs_path, exc)",
        "        _HANDLER_MOD_CACHE[ext] = None",
        "        return None",
        "",
        "",
        "def _pdf_read_text(path: str) -> str:",
        "    try:",
        "        import PyPDF2",
        "        reader = PyPDF2.PdfReader(path)",
        "        parts = []",
        "        for page in reader.pages:",
        "            try:",
        "                t = page.extract_text()",
        "                if t:",
        "                    parts.append(t)",
        "            except Exception:",
        "                pass",
        r'        return "\n".join(parts)',
        "    except Exception as exc:",
        '        logger.debug("PDF text extraction failed for %s: %s", path, exc)',
        '        return ""',
        "",
        "",
        "def _docx_read_text(path: str) -> str:",
        "    try:",
        "        from docx import Document",
        "        doc = Document(path)",
        r'        return "\n".join(p.text for p in doc.paragraphs if p.text)',
        "    except Exception as exc:",
        '        logger.debug("DOCX text extraction failed for %s: %s", path, exc)',
        '        return ""',
        "",
        "",
        "def _pptx_read_text(path: str) -> str:",
        "    try:",
        "        from pptx import Presentation",
        "        prs = Presentation(path)",
        "        parts = []",
        "        for slide in prs.slides:",
        "            for shape in slide.shapes:",
        '                if hasattr(shape, "text") and shape.text:',
        "                    parts.append(shape.text)",
        r'        return "\n".join(parts)',
        "    except Exception as exc:",
        '        logger.debug("PPTX text extraction failed for %s: %s", path, exc)',
        '        return ""',
        "",
        "",
        "def _plain_read_text(path: str) -> str:",
        "    try:",
        '        return Path(path).read_text(encoding="utf-8", errors="replace")',
        "    except Exception as exc:",
        '        logger.debug("Plain text read failed for %s: %s", path, exc)',
        '        return ""',
        "",
        "",
        "def read_text(path: str) -> str:",
        "    cfg = get_handler(path)",
        "    if cfg is None or not cfg.extract_text:",
        '        return ""',
        "    p = Path(path)",
        "    if not p.exists():",
        '        return ""',
        "    ext = p.suffix.lstrip('.').lower()",
        "    _mod = _load_installed(ext)",
        "    if _mod is not None and hasattr(_mod, 'read_text'):",
        "        try:",
        "            _text = _mod.read_text(path)",
        "            if _text:",
        "                return _text",
        "        except Exception as _exc:",
        '            logger.debug("Installed reader for %s failed on %s: %s", ext, path, _exc)',
        '    if cfg.category == "document-pdf":',
        "        return _pdf_read_text(path)",
        '    elif cfg.category == "document-office":',
        '        if ext in ("docx", "doc"):',
        "            return _docx_read_text(path)",
        '        elif ext in ("pptx", "ppt"):',
        "            return _pptx_read_text(path)",
        "        else:",
        '            return ""',
        '    elif cfg.category == "text-plain":',
        "        return _plain_read_text(path)",
        "    else:",
        '        return ""',
        "",
        "",
        "def _delegate_extract_images(src_path: Path, file_type: str) -> List[Path]:",
        "    try:",
        "        import extract_embedded_images",
        '        out_dir = src_path.parent / f"_extracted_{src_path.stem}"',
        "        return extract_embedded_images.extract(src_path, out_dir, file_type)",
        "    except Exception as exc:",
        '        logger.debug("extract_embedded_images failed for %s: %s", src_path, exc)',
        "        return []",
        "",
        "",
        "def extract_images(path: str) -> List[Path]:",
        "    cfg = get_handler(path)",
        "    if cfg is None or not cfg.extract_images:",
        "        return []",
        "    p = Path(path)",
        "    if not p.exists():",
        "        return []",
        "    ext = p.suffix.lstrip('.').lower()",
        "    _mod = _load_installed(ext)",
        "    if _mod is not None and hasattr(_mod, 'extract_images'):",
        "        try:",
        "            _imgs = _mod.extract_images(path, str(p.parent))",
        "            if _imgs:",
        "                return [Path(x) for x in _imgs]",
        "        except Exception as _exc:",
        '            logger.debug("Installed image extractor for %s failed on %s: %s", ext, path, _exc)',
        '    if cfg.category in ("image-raster", "image-vector"):',
        "        return [p]",
        '    elif cfg.category in ("document-pdf", "document-office"):',
        "        return _delegate_extract_images(p, ext)",
        "    else:",
        "        return []",
        "",
        "",
        "def handle(path: str) -> HandlerResult:",
        "    cfg = get_handler(path)",
        "    if cfg is None:",
        '        return HandlerResult(text=None, images=[], skipped=True, category="unknown")',
        "    try:",
        "        text = read_text(path) if cfg.extract_text else ''",
        "        images = extract_images(path) if cfg.extract_images else []",
        "    except Exception as exc:",
        '        logger.warning("handle() error processing %s: %s", path, exc)',
        "        text = ''",
        "        images = []",
        "    return HandlerResult(",
        "        text=text,",
        "        images=images,",
        "        skipped=False,",
        "        category=cfg.category,",
        "    )",
        "",
    ]
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def generate_from_dict(
    file_type_config: Dict[str, Any],
    out_path: Path,
) -> Path:
    """Generate file_type_handlers.py from a config dict.

    Args:
        file_type_config: The file_type_config sub-dict from config.json,
                          or {} to use all defaults.
        out_path:         Where to write the generated script.

    Returns:
        Path of the written file.
    """
    out_path = Path(out_path)
    categories, ext_map, installed_readers = _apply_config(file_type_config)
    source = _render_source(categories, ext_map, installed_readers)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(source, encoding="utf-8")
    return out_path


def generate(
    workdir: Path,
    out_path: Optional[Path] = None,
) -> Path:
    """Generate file_type_handlers.py for the given workdir.

    Reads file_type_config from <workdir>/.vault-bridge/config.json.
    On any read/parse failure, falls back to built-in defaults.

    Args:
        workdir:  Working directory. Config is expected at
                  <workdir>/.vault-bridge/config.json.
        out_path: Where to write the generated script. Defaults to
                  <workdir>/scripts/file_type_handlers.py.

    Returns:
        Path of the written file.
    """
    workdir = Path(workdir)
    if out_path is None:
        out_path = workdir / "scripts" / "file_type_handlers.py"

    file_type_config = _load_file_type_config(workdir)
    return generate_from_dict(file_type_config, out_path)


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import sys

    _workdir = Path(sys.argv[1]) if len(sys.argv) > 1 else Path.cwd()
    _out = generate(_workdir)
    print(f"Generated: {_out}")
