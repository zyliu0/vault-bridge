"""Extract embedded images from container documents.

Supports:
  - PDF: uses PyPDF2 page.images (3.x API) with manual /XObject traversal fallback
  - DOCX: python-docx related_parts with image/* content_type
  - PPTX: python-pptx slide shapes with MSO_SHAPE_TYPE.PICTURE
  - XLSX: stub — returns []
  - Image types (jpg/png/gif/webp/bmp/etc.): returns [src_path] unchanged
  - Anything else: returns []

Never raises — returns [] on any internal failure.

Python 3.9 compatible.
"""
import logging
from pathlib import Path
from typing import List

logger = logging.getLogger(__name__)

# File types considered direct images (no container extraction needed)
_IMAGE_TYPES = frozenset({
    "jpg", "jpeg", "png", "gif", "webp", "bmp", "tiff", "tif",
    "heic", "heif", "psd", "ai", "svg",
})

# Container types we can extract from
_CONTAINER_TYPES = frozenset({"pdf", "docx", "pptx"})


def extract(src_path: Path, out_dir: Path, file_type: str) -> List[Path]:
    """Extract embedded images from src_path into out_dir.

    Args:
        src_path: Path to the source document.
        out_dir: Directory to write extracted images into (created if needed).
        file_type: Lowercase file type hint: 'pdf', 'docx', 'pptx', 'xlsx',
                   or any image type ('jpg', 'png', etc.).

    Returns:
        List of Paths to extracted/returned image files.
        For image types: [src_path] (unchanged, no extraction).
        For container types: list of written image paths (may be empty).
        For xlsx or unknown: [].

    Never raises — returns [] on any internal failure.
    """
    ft = file_type.lower().lstrip(".")

    # Image passthrough — no extraction needed
    if ft in _IMAGE_TYPES:
        return [src_path]

    if ft not in _CONTAINER_TYPES and ft != "xlsx":
        return []

    # Ensure output directory exists
    try:
        out_dir.mkdir(parents=True, exist_ok=True)
    except Exception as exc:
        logger.warning("extract: could not create out_dir %s: %s", out_dir, exc)
        return []

    if ft == "pdf":
        return _extract_pdf(src_path, out_dir)
    elif ft == "docx":
        return _extract_docx(src_path, out_dir)
    elif ft == "pptx":
        return _extract_pptx(src_path, out_dir)
    elif ft == "xlsx":
        # v16.4.0 (AUDIT-1b). Pre-v16.4 this was an explicit early
        # ``return []``. The audit's representative XLSX had 40
        # inline images that all dropped silently. openpyxl exposes
        # ``ws._images`` with ``.image.ref`` pointing at the raw
        # bytes; we round-trip through PIL so the output is in a
        # format the downstream compressor understands.
        return _extract_xlsx(src_path, out_dir)

    return []


def _sniff_extension(data: bytes) -> str:
    """Return a file extension based on image magic bytes."""
    if data[:2] == b"\xff\xd8":
        return ".jpg"
    if data[:4] == b"\x89PNG":
        return ".png"
    if data[:4] == b"GIF8":
        return ".gif"
    if len(data) >= 12 and data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return ".webp"
    return ".bin"


def _unique_path(out_dir: Path, stem: str, ext: str) -> Path:
    """Return a unique path in out_dir, incrementing suffix on collision."""
    candidate = out_dir / f"{stem}{ext}"
    if not candidate.exists():
        return candidate
    counter = 1
    while True:
        candidate = out_dir / f"{stem}-{counter}{ext}"
        if not candidate.exists():
            return candidate
        counter += 1


_PDF_PAGE_RENDER_CAP = 6
_PDF_PAGE_RENDER_DPI = 150


def _extract_pdf(src_path: Path, out_dir: Path) -> List[Path]:
    """Extract images from a PDF file using PyPDF2.

    v16.6.0 (Batch A1): when neither ``page.images`` nor manual
    ``/XObject`` traversal yields any images, fall through to a
    PyMuPDF page-render path. This catches scanned/drawing PDFs that
    have no embedded image XObjects but DO have visible content (the
    body of every CD-phase TLS site plan, every faxed contract, every
    archived 1990s memo). Caps at ``_PDF_PAGE_RENDER_CAP`` pages so a
    200-page scan doesn't run for minutes.
    """
    try:
        import PyPDF2
    except ImportError:
        logger.warning("extract_pdf: PyPDF2 not available")
        return _extract_pdf_pages_fallback(src_path, out_dir)

    results: List[Path] = []
    src_stem = src_path.stem

    try:
        reader = PyPDF2.PdfReader(str(src_path))
    except Exception as exc:
        logger.warning("extract_pdf: could not open %s: %s", src_path, exc)
        return _extract_pdf_pages_fallback(src_path, out_dir)

    for page_idx, page in enumerate(reader.pages):
        # Try the 3.x page.images API first
        try:
            page_images = list(page.images)
        except Exception:
            page_images = []

        if page_images:
            for img_idx, img_file in enumerate(page_images):
                try:
                    data = img_file.data
                    ext = _sniff_extension(data)
                    stem = f"{src_stem}--page{page_idx}--img{img_idx}"
                    out_path = _unique_path(out_dir, stem, ext)
                    out_path.write_bytes(data)
                    results.append(out_path)
                except Exception as exc:
                    logger.warning("extract_pdf: error writing image: %s", exc)
            continue

        # Fallback: manual /XObject /Subtype /Image traversal
        try:
            resources = page.get("/Resources")
            if resources is None:
                continue
            xobjects = resources.get("/XObject")
            if xobjects is None:
                continue
            # Resolve indirect reference if needed
            if hasattr(xobjects, "get_object"):
                xobjects = xobjects.get_object()
            for name in list(xobjects.keys()):
                try:
                    obj = xobjects[name]
                    if hasattr(obj, "get_object"):
                        obj = obj.get_object()
                    if not hasattr(obj, "get"):
                        continue
                    subtype = obj.get("/Subtype")
                    if str(subtype) != "/Image":
                        continue
                    data = obj.get_data()
                    ext = _sniff_extension(data)
                    stem = f"{src_stem}--page{page_idx}--xobj{name.lstrip('/')}"
                    out_path = _unique_path(out_dir, stem, ext)
                    out_path.write_bytes(data)
                    results.append(out_path)
                except Exception as exc:
                    logger.warning("extract_pdf: xobject error: %s", exc)
        except Exception as exc:
            logger.warning("extract_pdf: page %d traversal error: %s", page_idx, exc)

    if results:
        return results

    # v16.6.0 (Batch A1): no embedded images at all — render the first
    # N pages as JPEGs so the writing LLM has something to see.
    return _extract_pdf_pages_fallback(src_path, out_dir)


def _extract_pdf_pages_fallback(src_path: Path, out_dir: Path) -> List[Path]:
    """Render PDF pages to JPEGs via PyMuPDF as a last-resort image path.

    v16.6.0 (Batch A1). Used when neither PyPDF2's ``page.images`` API
    nor manual ``/XObject`` traversal produced any image — common for
    scanned PDFs, drawing PDFs, and faxed memos that lack image
    XObjects but still carry visible content. Returns ``[]`` when
    PyMuPDF isn't installed or the document fails to open.
    """
    try:
        import fitz  # type: ignore  # PyMuPDF
    except ImportError:
        logger.debug("PDF page-render fallback: PyMuPDF not available")
        return []

    rendered: List[Path] = []
    src_stem = src_path.stem
    try:
        doc = fitz.open(str(src_path))
    except Exception as exc:
        logger.debug("PDF page-render fallback: could not open %s: %s", src_path, exc)
        return []

    try:
        n_pages = min(len(doc), _PDF_PAGE_RENDER_CAP)
        for page_idx in range(n_pages):
            try:
                page = doc[page_idx]
                pix = page.get_pixmap(dpi=_PDF_PAGE_RENDER_DPI)
                stem = f"{src_stem}--page{page_idx}--render"
                out_path = _unique_path(out_dir, stem, ".jpg")
                # Pixmap.tobytes("jpeg") returns valid JPEG bytes when
                # available; some PyMuPDF builds expose only PNG, in
                # which case fall through to writing PNG and let the
                # downstream compressor convert.
                try:
                    out_path.write_bytes(pix.tobytes("jpeg"))
                except Exception:
                    out_path = out_path.with_suffix(".png")
                    out_path.write_bytes(pix.tobytes("png"))
                rendered.append(out_path)
            except Exception as exc:
                logger.debug(
                    "PDF page-render fallback: page %d failed for %s: %s",
                    page_idx, src_path, exc,
                )
    finally:
        try:
            doc.close()
        except Exception:
            pass
    return rendered


def _extract_docx(src_path: Path, out_dir: Path) -> List[Path]:
    """Extract embedded images from a DOCX file using python-docx."""
    try:
        from docx import Document
    except ImportError:
        logger.warning("extract_docx: python-docx not available")
        return []

    results: List[Path] = []
    src_stem = src_path.stem

    try:
        doc = Document(str(src_path))
    except Exception as exc:
        logger.warning("extract_docx: could not open %s: %s", src_path, exc)
        return []

    idx = 0
    try:
        for rel_id, part in doc.part.related_parts.items():
            content_type = getattr(part, "content_type", "") or ""
            if not content_type.startswith("image/"):
                continue
            try:
                data = part.blob
                # Determine extension from content_type or magic bytes
                if content_type == "image/jpeg":
                    ext = ".jpg"
                elif content_type == "image/png":
                    ext = ".png"
                elif content_type == "image/gif":
                    ext = ".gif"
                elif content_type == "image/webp":
                    ext = ".webp"
                else:
                    ext = _sniff_extension(data)
                stem = f"{src_stem}--embedded-{idx}"
                out_path = _unique_path(out_dir, stem, ext)
                out_path.write_bytes(data)
                results.append(out_path)
                idx += 1
            except Exception as exc:
                logger.warning("extract_docx: error writing part %s: %s", rel_id, exc)
    except Exception as exc:
        logger.warning("extract_docx: error iterating related_parts: %s", exc)

    return results


def _extract_pptx(src_path: Path, out_dir: Path) -> List[Path]:
    """Extract embedded images from a PPTX file using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        logger.warning("extract_pptx: python-pptx not available")
        return []

    results: List[Path] = []
    src_stem = src_path.stem

    try:
        prs = Presentation(str(src_path))
    except Exception as exc:
        logger.warning("extract_pptx: could not open %s: %s", src_path, exc)
        return []

    idx = 0
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            try:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                data = shape.image.blob
                ext = _sniff_extension(data)
                stem = f"{src_stem}--slide{slide_idx}--img{idx}"
                out_path = _unique_path(out_dir, stem, ext)
                out_path.write_bytes(data)
                results.append(out_path)
                idx += 1
            except Exception as exc:
                logger.warning(
                    "extract_pptx: error extracting shape %s slide %d: %s",
                    getattr(shape, "name", "?"), slide_idx, exc
                )

    return results


def _extract_xlsx(src_path: Path, out_dir: Path) -> List[Path]:
    """Extract embedded images from an XLSX file using openpyxl.

    v16.4.0 (AUDIT-1b). XLSX files use OOXML's ``xl/media/`` part
    just like docx/pptx — openpyxl exposes them via
    ``Worksheet._images`` (each entry has ``.image.ref`` pointing
    at the underlying ``Image`` object whose ``.fp`` holds the raw
    bytes). We sniff the magic bytes to derive a reasonable
    extension so the downstream compressor doesn't have to guess.

    Falls back gracefully on any failure — empty list, no exception
    propagates.
    """
    try:
        import openpyxl  # type: ignore
        from openpyxl.drawing.image import Image as OpenpyxlImage  # type: ignore
    except ImportError:
        logger.warning("extract_xlsx: openpyxl not available")
        return []

    results: List[Path] = []
    src_stem = src_path.stem

    try:
        wb = openpyxl.load_workbook(str(src_path), data_only=True)
    except Exception as exc:
        logger.warning("extract_xlsx: could not open %s: %s", src_path, exc)
        return []

    idx = 0
    try:
        for sheet_idx, ws in enumerate(wb.worksheets):
            for img in getattr(ws, "_images", []) or []:
                try:
                    # openpyxl image.ref can be a PIL Image, a path,
                    # or a BytesIO depending on how the workbook was
                    # parsed. Try the documented attribute paths in order.
                    raw = None
                    inner = getattr(img, "_data", None)
                    if callable(inner):
                        try:
                            raw = inner()
                        except Exception:
                            raw = None
                    if raw is None:
                        ref = getattr(img, "ref", None) or getattr(img, "image", None)
                        if hasattr(ref, "read"):
                            try:
                                ref.seek(0)
                            except Exception:
                                pass
                            raw = ref.read()
                        elif isinstance(ref, (str, Path)):
                            try:
                                raw = Path(ref).read_bytes()
                            except Exception:
                                raw = None
                    if not raw:
                        continue
                    ext = _sniff_extension(raw)
                    stem = f"{src_stem}--sheet{sheet_idx}--img{idx}"
                    out_path = _unique_path(out_dir, stem, ext)
                    out_path.write_bytes(raw)
                    results.append(out_path)
                    idx += 1
                except Exception as exc:
                    logger.warning("extract_xlsx: error on sheet %d image %d: %s",
                                   sheet_idx, idx, exc)
    finally:
        try:
            wb.close()
        except Exception:
            pass

    return results
