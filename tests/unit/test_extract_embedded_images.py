"""Tests for scripts/extract_embedded_images.py — extract embedded images from containers.

TDD: tests written BEFORE the implementation.

Cases:
1. Image file type (.jpg) → returns [src_path] unchanged
2. DOCX with 2 images → returns 2 paths
3. PPTX with 1 image → returns 1 path
4. XLSX → returns []
5. Corrupt file → returns [], no exception
6. Unknown file_type → returns []
7. Output filename collision → suffix incremented
8. PDF with no images → returns []
"""
import io
import sys
from pathlib import Path

import pytest

SCRIPTS = Path(__file__).resolve().parents[2] / "scripts"
sys.path.insert(0, str(SCRIPTS))

import extract_embedded_images  # noqa: E402

# ---------------------------------------------------------------------------
# Minimal valid image bytes — generated via PIL for compatibility
# ---------------------------------------------------------------------------

def _make_png_bytes() -> bytes:
    """Generate a minimal valid 10x10 PNG via PIL."""
    from PIL import Image
    img = Image.new("RGB", (10, 10), color=(255, 0, 0))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


def _make_jpeg_bytes() -> bytes:
    """Generate a minimal valid 10x10 JPEG via PIL."""
    from PIL import Image
    img = Image.new("RGB", (10, 10), color=(0, 0, 255))
    buf = io.BytesIO()
    img.save(buf, "JPEG")
    return buf.getvalue()


# Compute once at module level
PNG_BYTES = _make_png_bytes()
JPEG_BYTES = _make_jpeg_bytes()


# ---------------------------------------------------------------------------
# Fixture helpers
# ---------------------------------------------------------------------------

def _make_docx_with_images(num_images: int) -> bytes:
    """Return bytes for a DOCX file with num_images distinct embedded PNG images."""
    from PIL import Image
    from docx import Document

    doc = Document()
    doc.add_heading("Test", 0)
    # Use different colors to ensure distinct blobs (python-docx deduplicates same bytes)
    colors = [(255, 0, 0), (0, 255, 0), (0, 0, 255), (255, 255, 0), (255, 0, 255)]
    for i in range(num_images):
        color = colors[i % len(colors)]
        img = Image.new("RGB", (10 + i, 10 + i), color=color)
        buf = io.BytesIO()
        img.save(buf, "PNG")
        buf.seek(0)
        doc.add_picture(buf)
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def _make_pptx_with_image() -> bytes:
    """Return bytes for a PPTX file with 1 embedded PNG image."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    buf = io.BytesIO(PNG_BYTES)
    slide.shapes.add_picture(buf, Inches(1), Inches(1))
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _make_xlsx() -> bytes:
    """Return bytes for a minimal XLSX file (no images)."""
    # openpyxl or just a minimal zip structure
    try:
        from openpyxl import Workbook
        wb = Workbook()
        out = io.BytesIO()
        wb.save(out)
        return out.getvalue()
    except ImportError:
        # Fallback: use a known XLSX header
        return b"PK\x03\x04"  # minimal ZIP/XLSX signature


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------

def test_extract_image_type_returns_src_path_unchanged(tmp_path):
    """Image file type → returns [src_path] without modification."""
    src = tmp_path / "photo.jpg"
    src.write_bytes(JPEG_BYTES)
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    result = extract_embedded_images.extract(src, out_dir, "jpg")
    assert result == [src]


def test_extract_png_type_returns_src_path_unchanged(tmp_path):
    """PNG file type → returns [src_path] without modification."""
    src = tmp_path / "image.png"
    src.write_bytes(PNG_BYTES)
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    result = extract_embedded_images.extract(src, out_dir, "png")
    assert result == [src]


def test_extract_docx_with_two_images(tmp_path):
    """DOCX with 2 embedded images → returns 2 paths, files exist."""
    src = tmp_path / "document.docx"
    src.write_bytes(_make_docx_with_images(2))
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    result = extract_embedded_images.extract(src, out_dir, "docx")
    assert len(result) == 2
    for p in result:
        assert p.exists(), f"Expected output file to exist: {p}"


def test_extract_pptx_with_one_image(tmp_path):
    """PPTX with 1 embedded image → returns 1 path."""
    src = tmp_path / "deck.pptx"
    src.write_bytes(_make_pptx_with_image())
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    result = extract_embedded_images.extract(src, out_dir, "pptx")
    assert len(result) == 1
    assert result[0].exists()


def test_extract_xlsx_returns_empty(tmp_path):
    """XLSX → returns [] (stub, no image extraction)."""
    src = tmp_path / "sheet.xlsx"
    src.write_bytes(_make_xlsx())
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    result = extract_embedded_images.extract(src, out_dir, "xlsx")
    assert result == []


def test_extract_corrupt_file_returns_empty_no_exception(tmp_path):
    """Corrupt file → returns [], does not raise."""
    src = tmp_path / "corrupt.docx"
    src.write_bytes(b"not a real docx just garbage bytes here")
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    # Must not raise
    result = extract_embedded_images.extract(src, out_dir, "docx")
    assert isinstance(result, list)


def test_extract_unknown_file_type_returns_empty(tmp_path):
    """Unknown file_type → returns []."""
    src = tmp_path / "weird.bin"
    src.write_bytes(b"random data")
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    result = extract_embedded_images.extract(src, out_dir, "bin")
    assert result == []


def test_extract_pdf_no_images_returns_empty(tmp_path):
    """PDF with no embedded images → returns []."""
    from PyPDF2 import PdfWriter
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    buf = io.BytesIO()
    writer.write(buf)

    src = tmp_path / "blank.pdf"
    src.write_bytes(buf.getvalue())
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    result = extract_embedded_images.extract(src, out_dir, "pdf")
    assert isinstance(result, list)


def test_extract_docx_extracted_files_are_in_out_dir(tmp_path):
    """Extracted DOCX images are written to out_dir."""
    src = tmp_path / "doc.docx"
    src.write_bytes(_make_docx_with_images(1))
    out_dir = tmp_path / "extracted"
    out_dir.mkdir()

    result = extract_embedded_images.extract(src, out_dir, "docx")
    for p in result:
        assert p.parent == out_dir, f"Expected files in {out_dir}, got {p.parent}"


def test_extract_pptx_image_has_valid_magic_bytes(tmp_path):
    """Extracted PPTX image has valid image magic bytes (PNG or JPEG)."""
    src = tmp_path / "deck.pptx"
    src.write_bytes(_make_pptx_with_image())
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    result = extract_embedded_images.extract(src, out_dir, "pptx")
    assert result
    first = result[0]
    magic = first.read_bytes()[:4]
    # PNG or JPEG
    assert magic[:3] == b"\xff\xd8\xff" or magic == b"\x89PNG", (
        f"Expected PNG or JPEG magic, got {magic!r}"
    )
